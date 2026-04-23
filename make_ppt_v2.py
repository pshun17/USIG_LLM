"""
USIG Methodology PPT v2
- 5 Component Agent boxes (AI agent design)
- Arrows converging into SIGNAL AGENT
- Portfolio construction + backtest results
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

IN = 914400

# ── Palette ───────────────────────────────────────────────────
C_HEADER    = RGBColor(0x0D, 0x1B, 0x2A)   # very dark navy
C_MID_BLUE  = RGBColor(0x1A, 0x4A, 0x8A)
C_CYAN      = RGBColor(0x00, 0xC8, 0xF0)   # neon cyan
C_CYAN_DIM  = RGBColor(0x00, 0x6A, 0x90)
C_GOLD      = RGBColor(0xE0, 0xB0, 0x30)
C_WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
C_LIGHT_BG  = RGBColor(0xE8, 0xF3, 0xFF)
C_DARK_TEXT = RGBColor(0x1A, 0x1A, 0x2E)
C_GREEN     = RGBColor(0x00, 0xCC, 0x88)
C_SIGNAL    = RGBColor(0x4A, 0x00, 0x8A)   # deep purple for Signal Agent
C_SIGNAL_LT = RGBColor(0x7A, 0x30, 0xC0)
C_SIGNAL_BG = RGBColor(0x1A, 0x00, 0x3A)

AGENT_COLORS = [
    (RGBColor(0x00, 0x6A, 0xC0), RGBColor(0x80, 0xBE, 0xFF)),  # Bond TR  — blue
    (RGBColor(0x7B, 0x2F, 0xBE), RGBColor(0xCC, 0xA0, 0xFF)),  # AI Macro — purple
    (RGBColor(0x00, 0x8A, 0x5A), RGBColor(0x70, 0xDD, 0xAA)),  # Eq Mom   — green
    (RGBColor(0xB8, 0x60, 0x00), RGBColor(0xFF, 0xCC, 0x70)),  # Eq Fund  — amber
    (RGBColor(0xBE, 0x2F, 0x5A), RGBColor(0xFF, 0x90, 0xB4)),  # Sentiment— pink
]

def add_shape(l, t, w, h, shape_type=1, fill=None, line=None, lw=1.0):
    s = slide.shapes.add_shape(shape_type, int(l*IN), int(t*IN), int(w*IN), int(h*IN))
    if fill: s.fill.solid(); s.fill.fore_color.rgb = fill
    else: s.fill.background()
    if line: s.line.color.rgb = line; s.line.width = Pt(lw)
    else: s.line.fill.background()
    return s

def add_oval(l, t, w, h, fill=None, line=None, lw=1.0):
    return add_shape(l, t, w, h, shape_type=9, fill=fill, line=line, lw=lw)

def add_hex(l, t, w, h, fill=None, line=None, lw=1.5):
    return add_shape(l, t, w, h, shape_type=125, fill=fill, line=line, lw=lw)

def add_text(l, t, w, h, text, size, bold=False, color=C_WHITE,
             align=PP_ALIGN.CENTER, italic=False):
    txb = slide.shapes.add_textbox(int(l*IN), int(t*IN), int(w*IN), int(h*IN))
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = "Calibri"
    return txb

def add_multiline(l, t, w, h, lines, size, bold=False, color=C_WHITE,
                  align=PP_ALIGN.LEFT, spacing=1.15):
    from pptx.util import Pt as _Pt
    from pptx.oxml.ns import qn
    from lxml import etree
    txb = slide.shapes.add_textbox(int(l*IN), int(t*IN), int(w*IN), int(h*IN))
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.size = _Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.name = "Calibri"
    return txb

# ═══════════════════════════════════════════════════════════════
# 1. BACKGROUND gradient blocks
# ═══════════════════════════════════════════════════════════════
add_shape(0, 0, 13.33, 7.5, fill=RGBColor(0xF0, 0xF5, 0xFF))   # base
add_shape(0, 0, 13.33, 0.78, fill=C_HEADER)                     # header
add_shape(0, 0.78, 13.33, 0.34, fill=C_MID_BLUE)                # sub-header
add_shape(0, 7.28, 13.33, 0.22, fill=C_HEADER)                  # footer

# ═══════════════════════════════════════════════════════════════
# 2. HEADER
# ═══════════════════════════════════════════════════════════════
add_text(0.28, 0.10, 10, 0.45,
         "USIG Credit Portfolio  |  Integrated Scoring Methodology",
         19, bold=True, color=C_WHITE, align=PP_ALIGN.LEFT)
add_text(10.20, 0.16, 2.90, 0.32,
         "April 2026", 10, color=RGBColor(0x90, 0xBC, 0xE8), align=PP_ALIGN.RIGHT)

# Gold accent line
add_shape(0, 0.78, 13.33, 0.04, fill=C_GOLD)

add_text(0.28, 0.82, 12.77, 0.28,
         "Integrated Score  =  Σ  Component Agent Scores  →  Signal Agent  →  Portfolio Selection",
         10.5, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════
# 3. COMPONENT AGENT BOXES
# ═══════════════════════════════════════════════════════════════
AGENTS = [
    ("Bond TR\nScore",  "W₁ = 50%",
     ["• Carry (2.5M YTW)", "• OAS Compression", "• DP Rating Spread", "→ Class %ile rank"],
     "01"),
    ("AI Macro\nScore", "W₂ = 50%",
     ["• Sector view ×0.40", "• Maturity pref ×0.35", "• Rating buffer ×0.25", "→ Clipped [−1,+1]"],
     "02"),
    ("Eq Mom\nScore",   "W₃  opt.",
     ["• 1M/3M return", "• 30D volatility", "• 52W high ratio", "→ Cross-sect. rank"],
     "03"),
    ("Eq Fund\nScore",  "W₄  opt.",
     ["• D/E · Profit Margin", "• Revenue Growth", "• CR · EV/EBITDA", "→ Cross-sect. rank"],
     "04"),
    ("Sentiment\nScore","W₅  opt.",
     ["• News VADER score", "• Google Trends", "→ Cross-sect. rank", "⚠ Backtest: N/A"],
     "05"),
]

BOX_Y  = 1.22
BOX_W  = 2.38
BOX_H  = 1.96
GAP    = 0.18
START_X = (13.33 - (5*BOX_W + 4*GAP)) / 2

box_centers = []
for i, (title, weight, bullets, num) in enumerate(AGENTS):
    col_dark, col_light = AGENT_COLORS[i]
    x = START_X + i*(BOX_W + GAP)
    cx = x + BOX_W/2
    box_centers.append(cx)

    # Shadow
    add_shape(x+0.06, BOX_Y+0.06, BOX_W, BOX_H,
              fill=RGBColor(0xC0, 0xD0, 0xE8))
    # Box body
    add_shape(x, BOX_Y, BOX_W, BOX_H,
              fill=RGBColor(0xF5, 0xF9, 0xFF), line=col_dark, lw=1.2)
    # Header strip
    add_shape(x, BOX_Y, BOX_W, 0.46, fill=col_dark)
    # Accent line below header
    add_shape(x, BOX_Y+0.46, BOX_W, 0.03, fill=col_light)
    # Title
    add_text(x, BOX_Y+0.02, BOX_W, 0.42, title,
             10, bold=True, color=C_WHITE)
    # Weight tag
    add_shape(x + BOX_W/2 - 0.45, BOX_Y+0.52, 0.90, 0.20,
              fill=col_dark, line=col_light, lw=0.6)
    add_text(x + BOX_W/2 - 0.45, BOX_Y+0.52, 0.90, 0.20,
             weight, 7.5, bold=True, color=col_light)
    # Bullets
    add_multiline(x+0.10, BOX_Y+0.76, BOX_W-0.15, BOX_H-0.80,
                  bullets, 7.8, color=C_DARK_TEXT, align=PP_ALIGN.LEFT)
    # Bottom circuit line
    add_shape(x+0.08, BOX_Y+BOX_H-0.07, BOX_W-0.16, 0.03, fill=col_light)
    # Corner nodes
    for dx in [0.05, BOX_W-0.14]:
        add_oval(x+dx, BOX_Y+BOX_H-0.10, 0.09, 0.09, fill=col_light)

    # ── Hexagon agent node (above box) ──────────────────────
    HEX_D = 0.36
    hx = cx - HEX_D/2
    hy = BOX_Y - HEX_D - 0.12
    # Glow ring
    add_oval(hx-0.04, hy-0.04, HEX_D+0.08, HEX_D+0.08, fill=col_light)
    # Hex
    add_hex(hx, hy, HEX_D, HEX_D, fill=col_dark, line=C_WHITE, lw=0.8)
    add_text(hx, hy+0.06, HEX_D, HEX_D-0.12, num, 8, bold=True, color=C_WHITE)
    # Vertical line node → box
    add_shape(cx-0.015, hy+HEX_D, 0.03, BOX_Y-(hy+HEX_D), fill=col_dark)
    # Arrowhead
    add_shape(cx-0.05, BOX_Y-0.07, 0.10, 0.07, fill=col_dark)

    # ── AGENT badge (bottom-right inside box) ───────────────
    add_shape(x+BOX_W-0.62, BOX_Y+BOX_H-0.30, 0.54, 0.20,
              fill=col_dark, line=col_light, lw=0.5)
    add_text(x+BOX_W-0.62, BOX_Y+BOX_H-0.30, 0.54, 0.20,
             f"AGENT {num}", 6, bold=True, color=col_light)

# Pipeline spine connecting all hex nodes
PIPE_Y = BOX_Y - HEX_D/2 - 0.12 + 0.18   # center of hex nodes row
add_shape(START_X-0.10, PIPE_Y-0.015, 13.33-START_X*2+0.20, 0.03,
          fill=C_CYAN_DIM)
add_shape(START_X-0.10, PIPE_Y-0.007, 13.33-START_X*2+0.20, 0.014,
          fill=C_CYAN)
add_shape(0.05, PIPE_Y-0.17, 1.08, 0.34, fill=C_HEADER)
add_shape(0.05, PIPE_Y-0.17, 1.08, 0.34, fill=None,
          line=C_CYAN, lw=0.7)
add_text(0.05, PIPE_Y-0.14, 1.08, 0.28,
         "SCORING\nPIPELINE", 6, bold=True, color=C_CYAN)

# Chevrons between nodes
for i in range(len(box_centers)-1):
    mx = (box_centers[i] + box_centers[i+1])/2 - 0.05
    add_text(mx, PIPE_Y-0.14, 0.10, 0.28, "▶", 7,
             color=C_HEADER, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════
# 4. SIGNAL AGENT (central aggregator)
# ═══════════════════════════════════════════════════════════════
SA_CX = 13.33/2
SA_W  = 3.80
SA_H  = 0.80
SA_Y  = BOX_Y + BOX_H + 0.48
SA_X  = SA_CX - SA_W/2

# Converging arrows from each box to signal agent
for i, cx in enumerate(box_centers):
    # Dashed-style thin line from box bottom to signal agent top
    by = BOX_Y + BOX_H + 0.06
    ay = SA_Y - 0.06
    # Vertical segment down from box
    seg_h = (SA_Y - 0.06) - (BOX_Y + BOX_H + 0.06)
    add_shape(cx-0.015, by, 0.03, seg_h*0.5,
              fill=AGENT_COLORS[i][0])
    # Diagonal to center: use horizontal then vertical
    mid_y = by + seg_h*0.5
    # Horizontal toward center
    direction = SA_CX - cx
    add_shape(min(cx, SA_CX)-0.015, mid_y-0.015,
              abs(direction)+0.03, 0.03,
              fill=AGENT_COLORS[i][0])
    # Vertical down to signal agent
    add_shape(SA_CX-0.015, mid_y, 0.03, SA_Y - mid_y - 0.04,
              fill=C_SIGNAL_LT)

# Arrowhead into Signal Agent
add_shape(SA_CX-0.06, SA_Y-0.09, 0.12, 0.09, fill=C_SIGNAL_LT)

# Signal Agent glow
add_shape(SA_X-0.06, SA_Y-0.06, SA_W+0.12, SA_H+0.12,
          fill=RGBColor(0xCC, 0xAA, 0xFF), line=None)
# Signal Agent box
add_shape(SA_X, SA_Y, SA_W, SA_H, fill=C_SIGNAL_BG, line=C_CYAN, lw=2.0)
# Inner highlight top bar
add_shape(SA_X, SA_Y, SA_W, 0.08, fill=C_CYAN)
# Label
add_text(SA_X, SA_Y+0.06, SA_W, 0.30,
         "⚡  SIGNAL AGENT", 14, bold=True, color=C_CYAN)
add_text(SA_X, SA_Y+0.38, SA_W, 0.30,
         "Weighted Aggregation  →  Integrated Score  →  Top-N Selection",
         8.5, bold=False, color=RGBColor(0xCC, 0xCC, 0xFF))
# Corner hex decorations
for hx_off in [0.08, SA_W-0.36]:
    add_hex(SA_X+hx_off, SA_Y+0.18, 0.28, 0.28,
            fill=C_SIGNAL_LT, line=C_CYAN, lw=0.8)

# ═══════════════════════════════════════════════════════════════
# 5. OUTPUT ARROW + PORTFOLIO CONSTRUCTION
# ═══════════════════════════════════════════════════════════════
OUT_Y = SA_Y + SA_H + 0.10
add_shape(SA_CX-0.015, OUT_Y, 0.03, 0.22, fill=C_CYAN)
add_shape(SA_CX-0.07, OUT_Y+0.18, 0.14, 0.10, fill=C_CYAN)

PROC_Y = OUT_Y + 0.32
PROC_H = 0.82
PROC_W = 3.80
procs = [
    (0.35,  "① CLASS Definition",
     "BCLASS3 × Maturity bucket × Rating tier"),
    (4.77,  "② Bond Selection",
     "Top 3 per CLASS by Integrated Score\nEqual-weight · Monthly rebalance"),
    (9.18,  "③ vs Benchmark",
     "Mkt-val weighted Total Return\n(full LUACSTAT universe)"),
]
for px, ptitle, pbody in procs:
    add_shape(px, PROC_Y, PROC_W, PROC_H,
              fill=C_HEADER, line=C_GOLD, lw=1.2)
    add_text(px+0.12, PROC_Y+0.05, PROC_W-0.20, 0.26,
             ptitle, 9.5, bold=True, color=C_GOLD)
    add_text(px+0.12, PROC_Y+0.32, PROC_W-0.20, 0.44,
             pbody, 8, color=C_WHITE)

# Arrows between proc boxes
for arr_x in [4.17+0.20, 8.58+0.20]:
    add_text(arr_x, PROC_Y+0.28, 0.20, 0.26, "▶", 11,
             color=C_CYAN, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════
# 6. BACKTEST RESULT STRIP
# ═══════════════════════════════════════════════════════════════
RES_Y = PROC_Y + PROC_H + 0.12
RES_H = 0.68
add_shape(0.35, RES_Y, 12.63, RES_H,
          fill=RGBColor(0xF8, 0xFB, 0xFF), line=C_MID_BLUE, lw=0.8)
add_shape(0.35, RES_Y, 12.63, 0.06, fill=C_MID_BLUE)
add_text(0.35, RES_Y+0.06, 12.63, 0.22,
         "Backtest Result  (Bond TR + AI Macro 50/50)  |  Jan 2021 – Mar 2026  (63 months)",
         9, bold=True, color=C_MID_BLUE, align=PP_ALIGN.CENTER)
metrics = [
    ("Cumulative Return", "Model +5.9%  vs  Bmk +4.7%"),
    ("Ann. Alpha",        "+1.17% / yr"),
    ("Information Ratio", "1.11"),
    ("Hit Rate",          "59%  (monthly)"),
]
for i, (k, v) in enumerate(metrics):
    mx = 0.55 + i*3.10
    add_text(mx, RES_Y+0.30, 2.90, 0.18, k,
             7.5, color=RGBColor(0x55,0x55,0x88), align=PP_ALIGN.CENTER)
    add_text(mx, RES_Y+0.46, 2.90, 0.20, v,
             9.5, bold=True, color=C_MID_BLUE, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════
# 7. FOOTER
# ═══════════════════════════════════════════════════════════════
add_text(0, 7.29, 13.33, 0.20,
         "USIG Credit Research  |  Confidential  |  For internal use only",
         7.5, color=RGBColor(0x88, 0xAA, 0xCC), align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════
out = r"C:\Users\sh.park\Documents\USIG_LLM\methodology_overview_v2.pptx"
prs.save(out)
print(f"Saved → {out}")
