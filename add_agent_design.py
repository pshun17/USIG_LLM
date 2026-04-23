"""
Adds AI-agent visual design to the 5 component boxes in methodology_overview.pptx.
Opens existing file, adds elements on top, saves as methodology_overview_v2.pptx.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree
import copy

# ── Colors ────────────────────────────────────────────────────
CYAN        = RGBColor(0x00, 0xD4, 0xFF)   # neon cyan — pipeline / nodes
DARK_BLUE   = RGBColor(0x1F, 0x3A, 0x5F)
MID_BLUE    = RGBColor(0x2E, 0x5E, 0xA8)
NEON_GREEN  = RGBColor(0x00, 0xE5, 0x96)   # output node
GOLD        = RGBColor(0xC9, 0xA0, 0x2A)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
DARK_BG     = RGBColor(0x0A, 0x14, 0x2A)   # deep navy for node bg
LIGHT_CYAN  = RGBColor(0xD0, 0xF5, 0xFF)

IN = 914400   # EMUs per inch

prs = Presentation(r"C:\Users\sh.park\Documents\USIG_LLM\methodology_overview.pptx")
slide = prs.slides[0]

# ── helper: add a filled+bordered shape ──────────────────────
def rect(l, t, w, h, fill=None, line=None, lw=1.0, rounding=None):
    if rounding:
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        shp = slide.shapes.add_shape(5, int(l*IN), int(t*IN), int(w*IN), int(h*IN))  # rounded rect
        shp.adjustments[0] = rounding
    else:
        shp = slide.shapes.add_shape(1, int(l*IN), int(t*IN), int(w*IN), int(h*IN))
    if fill:
        shp.fill.solid(); shp.fill.fore_color.rgb = fill
    else:
        shp.fill.background()
    if line:
        shp.line.color.rgb = line; shp.line.width = Pt(lw)
    else:
        shp.line.fill.background()
    return shp

def oval(l, t, w, h, fill=None, line=None, lw=1.0):
    shp = slide.shapes.add_shape(9, int(l*IN), int(t*IN), int(w*IN), int(h*IN))  # oval
    if fill:
        shp.fill.solid(); shp.fill.fore_color.rgb = fill
    else:
        shp.fill.background()
    if line:
        shp.line.color.rgb = line; shp.line.width = Pt(lw)
    else:
        shp.line.fill.background()
    return shp

def textbox(l, t, w, h, text, size, bold=False, color=WHITE, align=PP_ALIGN.CENTER):
    txb = slide.shapes.add_textbox(int(l*IN), int(t*IN), int(w*IN), int(h*IN))
    txb.word_wrap = True
    p = txb.text_frame.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Calibri"
    return txb

def hex_shape(l, t, w, h, fill=None, line=None, lw=1.5):
    """Hexagon (shape type 125)"""
    shp = slide.shapes.add_shape(125, int(l*IN), int(t*IN), int(w*IN), int(h*IN))
    if fill:
        shp.fill.solid(); shp.fill.fore_color.rgb = fill
    else:
        shp.fill.background()
    if line:
        shp.line.color.rgb = line; shp.line.width = Pt(lw)
    else:
        shp.line.fill.background()
    return shp

# ── Box layout (from shape analysis) ─────────────────────────
boxes = [
    # (x,     w,    label,            icon)
    (0.15,  2.45, "Bond TR",         "📊"),
    (2.72,  2.45, "AI Macro",        "🤖"),
    (5.29,  2.45, "Eq Mom",          "📈"),
    (7.86,  2.45, "Eq Fund",         "🏦"),
    (10.43, 2.75, "Sentiment",       "💬"),
]
BOX_TOP = 1.40

centers = [x + w/2 for x, w, *_ in boxes]

# ═══════════════════════════════════════════════════════════════
# A. PIPELINE SPINE — thin horizontal neon line above all boxes
# ═══════════════════════════════════════════════════════════════
PIPE_Y    = 0.98   # center of pipeline
NODE_D    = 0.38   # node circle diameter
NODE_TOP  = PIPE_Y - NODE_D/2

# Background spine bar (dark navy)
rect(0.05, PIPE_Y - 0.03, 13.23, 0.06, fill=DARK_BG)
# Neon cyan spine line
rect(0.05, PIPE_Y - 0.015, 13.23, 0.03, fill=CYAN)

# "AGENT PIPELINE" label on left
rect(0.05, PIPE_Y - 0.20, 1.40, 0.40, fill=DARK_BG, line=CYAN, lw=0.8)
textbox(0.05, PIPE_Y - 0.17, 1.40, 0.34, "AGENT  PIPELINE", 6.5,
        bold=True, color=CYAN, align=PP_ALIGN.CENTER)

# "→ Integrated Score" node on right
right_x = 11.95
oval(right_x, NODE_TOP, NODE_D*1.4, NODE_D, fill=NEON_GREEN, line=WHITE, lw=0.8)
textbox(right_x, NODE_TOP + 0.04, NODE_D*1.4, NODE_D - 0.08,
        "∑ Score", 6.5, bold=True, color=DARK_BG, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════
# B. AGENT NODES — hexagon per box, sitting on the pipeline
# ═══════════════════════════════════════════════════════════════
agent_labels = ["BT", "AI", "EM", "EF", "ST"]
agent_colors = [
    RGBColor(0x00, 0x8B, 0xD4),   # Bond TR — steel blue
    RGBColor(0x7B, 0x2F, 0xBE),   # AI Macro — purple
    RGBColor(0x00, 0xBF, 0x7A),   # Eq Mom   — green
    RGBColor(0xC0, 0x6B, 0x00),   # Eq Fund  — amber
    RGBColor(0xBE, 0x2F, 0x5A),   # Sentiment — red
]

for i, (cx, (x, w, label, icon)) in enumerate(zip(centers, boxes)):
    col = agent_colors[i]
    node_l = cx - NODE_D/2

    # Outer glow ring (larger, semi-transparent via color)
    glow_col = RGBColor(
        min(col.red   + 60, 255),
        min(col.green + 60, 255),
        min(col.blue  + 60, 255)
    )
    oval(node_l - 0.04, NODE_TOP - 0.04,
         NODE_D + 0.08, NODE_D + 0.08,
         fill=glow_col, line=None)

    # Hexagon node
    hex_shape(node_l, NODE_TOP, NODE_D, NODE_D, fill=col, line=WHITE, lw=0.8)

    # Agent label inside hexagon
    textbox(node_l, NODE_TOP + 0.06, NODE_D, NODE_D - 0.12,
            agent_labels[i], 8.5, bold=True, color=WHITE)

    # Vertical connector from node bottom → box top
    vline_x = cx - 0.015
    vline_top = NODE_TOP + NODE_D
    vline_h   = BOX_TOP - vline_top
    if vline_h > 0:
        rect(vline_x, vline_top, 0.03, vline_h, fill=col)
        # Small triangle arrowhead at bottom
        rect(vline_x - 0.04, BOX_TOP - 0.08, 0.11, 0.08, fill=col)

    # ── C. AGENT BADGE inside box (bottom-right corner) ──────
    badge_x = x + w - 0.60
    badge_y = BOX_TOP + 2.15 - 0.32   # near bottom of box
    rect(badge_x, badge_y, 0.52, 0.24,
         fill=col, line=None, rounding=0.05)
    textbox(badge_x, badge_y + 0.01, 0.52, 0.22,
            f"AGENT {i+1:02d}", 6.0, bold=True, color=WHITE)

    # ── D. Circuit-style top border accent ───────────────────
    # Thin colored line just below the box header (y = BOX_TOP + 0.38)
    accent_y = BOX_TOP + 0.38
    rect(x + 0.05, accent_y, w - 0.10, 0.03, fill=col)

    # Small corner dots (circuit nodes) — top-left, top-right
    dot_d = 0.07
    oval(x + 0.05,           accent_y - dot_d/2, dot_d, dot_d, fill=col)
    oval(x + w - 0.05 - dot_d, accent_y - dot_d/2, dot_d, dot_d, fill=col)

# ═══════════════════════════════════════════════════════════════
# E. Between-box connectors (dashed flow arrows on pipeline)
# ═══════════════════════════════════════════════════════════════
for i in range(len(centers) - 1):
    mid_x = (centers[i] + centers[i+1]) / 2 - 0.06
    # Small chevron/arrow
    textbox(mid_x, PIPE_Y - 0.14, 0.12, 0.28,
            "▶", 7, bold=False, color=DARK_BG, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════
# F. "AI AGENT SCORING ENGINE" watermark label
# ═══════════════════════════════════════════════════════════════
rect(0.15, 1.40, 13.03, 2.15, fill=None,
     line=RGBColor(0x00, 0xD4, 0xFF), lw=0.6)

# Top-left label on the outer border
rect(0.15, 1.32, 2.20, 0.20, fill=DARK_BG)
textbox(0.17, 1.32, 2.16, 0.20,
        "AI AGENT SCORING ENGINE", 6.0, bold=True, color=CYAN)

# ═══════════════════════════════════════════════════════════════
out = r"C:\Users\sh.park\Documents\USIG_LLM\methodology_overview_v2.pptx"
prs.save(out)
print(f"Saved → {out}")
