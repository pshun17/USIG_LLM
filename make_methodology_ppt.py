from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import pptx.oxml.ns as nsmap
from lxml import etree

# ── Color palette ──────────────────────────────────────────────
DARK_BLUE   = RGBColor(0x1F, 0x3A, 0x5F)   # header bg
MID_BLUE    = RGBColor(0x2E, 0x5E, 0xA8)   # section titles
LIGHT_BLUE  = RGBColor(0xD6, 0xE4, 0xF7)   # box bg
GOLD        = RGBColor(0xC9, 0xA0, 0x2A)   # accent
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
DARK_GRAY   = RGBColor(0x2B, 0x2B, 0x2B)
MID_GRAY    = RGBColor(0x55, 0x55, 0x55)
BOX_BORDER  = RGBColor(0x2E, 0x5E, 0xA8)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

slide_layout = prs.slide_layouts[6]   # blank
slide = prs.slides.add_slide(slide_layout)

def add_rect(slide, l, t, w, h, fill_rgb=None, line_rgb=None, line_width_pt=1.0):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    if fill_rgb:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_rgb
    else:
        shape.fill.background()
    if line_rgb:
        shape.line.color.rgb = line_rgb
        shape.line.width = Pt(line_width_pt)
    else:
        shape.line.fill.background()
    return shape

def add_textbox(slide, l, t, w, h, text, size, bold=False, color=DARK_GRAY,
                align=PP_ALIGN.LEFT, wrap=True):
    txBox = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txBox.word_wrap = wrap
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Calibri"
    return txBox

def add_text_in_shape(shape, text, size, bold=False, color=WHITE, align=PP_ALIGN.CENTER):
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Calibri"

# ── 1. Header bar ──────────────────────────────────────────────
hdr = add_rect(slide, 0, 0, 13.33, 0.85, fill_rgb=DARK_BLUE)
add_textbox(slide, 0.25, 0.08, 9, 0.55,
            "USIG Credit Portfolio | Integrated Scoring Methodology",
            20, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
add_textbox(slide, 10.0, 0.15, 3.1, 0.45,
            "As of April 2026", 11, bold=False,
            color=RGBColor(0xAA, 0xC8, 0xF0), align=PP_ALIGN.RIGHT)

# ── 2. Sub-header: formula ─────────────────────────────────────
add_rect(slide, 0, 0.85, 13.33, 0.42, fill_rgb=MID_BLUE)
add_textbox(slide, 0.25, 0.90, 12.8, 0.35,
            "Integrated Score  =  Bond TR Score × W₁  +  AI Macro Score × W₂  "
            "(+  Eq Mom × W₃  +  Eq Fund × W₄  +  Sentiment × W₅  when data available)",
            11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# ── 3. Score component boxes ───────────────────────────────────
# positions: 5 boxes across, y=1.40, height=2.10
BOX_Y = 1.40
BOX_H = 2.15
box_cfg = [
    # (x,    w,   label,        weight, items)
    (0.15, 2.45, "Bond TR Score",   "W₁ = 50%",
     ["Carry (2.5M YTW)",
      "OAS Compression",
      "DP Rating Spread",
      "→ Class-level %ile rank"]),
    (2.72, 2.45, "AI Macro Score",  "W₂ = 50%",
     ["Sector view × 0.40",
      "Maturity preference × 0.35",
      "Rating buffer × 0.25",
      "→ Clipped [ −1, +1 ]"]),
    (5.29, 2.45, "Eq Mom Score",    "W₃  (optional)",
     ["1M / 3M price return",
      "30-day volatility",
      "52-week high ratio",
      "→ Cross-sectional rank"]),
    (7.86, 2.45, "Eq Fund Score",   "W₄  (optional)",
     ["D/E · Profit Margin",
      "Revenue Growth",
      "Current Ratio · EV/EBITDA",
      "→ Cross-sectional rank"]),
    (10.43, 2.75, "Sentiment Score", "W₅  (optional)",
     ["News headline VADER",
      "Google Trends signals",
      "→ Cross-sectional rank",
      "⚠ Backtest: not available"]),
]

for (x, w, label, weight, items) in box_cfg:
    # box background
    box = add_rect(slide, x, BOX_Y, w, BOX_H,
                   fill_rgb=LIGHT_BLUE, line_rgb=BOX_BORDER, line_width_pt=1.2)
    # label header inside box
    lbl_box = add_rect(slide, x, BOX_Y, w, 0.38, fill_rgb=MID_BLUE)
    add_text_in_shape(lbl_box, label, 10.5, bold=True, color=WHITE)
    # weight tag
    add_textbox(slide, x, BOX_Y + 0.40, w, 0.28,
                weight, 9.5, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    # bullet items
    for i, item in enumerate(items):
        add_textbox(slide, x + 0.10, BOX_Y + 0.72 + i * 0.33, w - 0.15, 0.32,
                    "• " + item, 8.5, bold=False, color=DARK_GRAY)

# ── 4. Arrow → Portfolio Selection ────────────────────────────
# Arrow label
add_textbox(slide, 0.15, 3.72, 13.0, 0.30,
            "▼  Portfolio Construction", 10, bold=True,
            color=MID_BLUE, align=PP_ALIGN.CENTER)

# ── 5. Bottom row: 3 process boxes ────────────────────────────
PROC_Y = 4.08
PROC_H = 1.10
proc_boxes = [
    (0.30, 3.90, "① CLASS Definition",
     "BCLASS3 × Maturity bucket × Rating tier\n(~3Y / 3–7Y / 7–15Y / 15Y+  |  A↑ vs BBB↓)"),
    (4.55, 3.90, "② Bond Selection",
     "Top 3 bonds per CLASS by Integrated Score\nEqual-weighted within class · Monthly rebalance"),
    (8.80, 3.90, "③ Benchmark",
     "Mkt-value weighted avg Total Return\nof entire LUACSTAT universe (all bonds)"),
]
for (x, w, title, body) in proc_boxes:
    pb = add_rect(slide, x, PROC_Y, w, PROC_H,
                  fill_rgb=DARK_BLUE, line_rgb=GOLD, line_width_pt=1.5)
    add_textbox(slide, x + 0.12, PROC_Y + 0.05, w - 0.20, 0.30,
                title, 10, bold=True, color=GOLD)
    add_textbox(slide, x + 0.12, PROC_Y + 0.38, w - 0.20, 0.65,
                body, 8.5, bold=False, color=WHITE)

# ── 6. Backtest result banner ──────────────────────────────────
res_y = 5.35
add_rect(slide, 0.30, res_y, 12.73, 0.78,
         fill_rgb=RGBColor(0xF0, 0xF5, 0xFF), line_rgb=MID_BLUE, line_width_pt=1.0)

add_textbox(slide, 0.45, res_y + 0.05, 12.4, 0.28,
            "Backtest Result  (Bond TR + AI Macro, 50/50)  |  Jan 2021 – Mar 2026  (63 months)",
            9.5, bold=True, color=MID_BLUE, align=PP_ALIGN.CENTER)

metrics = [
    ("Cumulative Return", "Model  +5.9%   vs   Bmk  +4.7%"),
    ("Ann. Alpha",        "+1.17% / yr"),
    ("Information Ratio", "1.11"),
    ("Hit Rate",          "59%  (monthly)"),
]
col_w = 3.0
for i, (k, v) in enumerate(metrics):
    cx = 0.55 + i * col_w
    add_textbox(slide, cx, res_y + 0.35, col_w - 0.1, 0.22,
                k, 7.5, bold=False, color=MID_GRAY, align=PP_ALIGN.CENTER)
    add_textbox(slide, cx, res_y + 0.52, col_w - 0.1, 0.24,
                v, 9.5, bold=True, color=DARK_BLUE, align=PP_ALIGN.CENTER)

# ── 7. Footer ──────────────────────────────────────────────────
add_rect(slide, 0, 6.25, 13.33, 0.25, fill_rgb=DARK_BLUE)
add_textbox(slide, 0.25, 6.27, 12.8, 0.20,
            "USIG Credit Research  |  Confidential  |  For internal use only",
            7.5, bold=False, color=RGBColor(0xAA, 0xC4, 0xE8),
            align=PP_ALIGN.CENTER)

out = r"C:\Users\sh.park\Documents\USIG_LLM\methodology_overview.pptx"
prs.save(out)
print(f"Saved: {out}")
